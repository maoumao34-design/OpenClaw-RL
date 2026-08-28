# -*- coding: utf-8 -*-
"""只更新第8页：清空正文，插入 harvest_flow.png（方案 B 流程图）。

不动其它页。可重复运行。

用法：python fill_slide8.py
"""

import shutil
from datetime import datetime

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

FILE = "实习小结_毛泽辉.pptx"
IMG = "harvest_flow.png"

EA_FONT = "微软雅黑"
LATIN_FONT = "Segoe UI"

# 标题下方整幅内容区；图长宽比 1952/880 ≈ 2.218
IMG_LEFT, IMG_TOP = 0.55, 1.55
IMG_WIDTH = 12.20


def style_run(run, size, bold=None):
    font = run.font
    font.size = size
    if bold is not None:
        font.bold = bold
    font.name = LATIN_FONT
    rPr = font._rPr
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        latin = rPr.find(qn("a:latin"))
        latin.addnext(ea) if latin is not None else rPr.append(ea)
    ea.set("typeface", EA_FONT)
    rPr.set("lang", "zh-CN")
    rPr.set("altLang", "en-US")
    rPr.set("noProof", "1")


def disable_proofing(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rPr = run.font._rPr
                rPr.set("lang", "zh-CN")
                rPr.set("altLang", "en-US")
                rPr.set("noProof", "1")


def main():
    # 改前备份
    stamp = datetime.now().strftime("%H%M")
    bak = f"实习小结_毛泽辉_backup_改第8页前_{stamp}.pptx"
    shutil.copy2(FILE, bak)
    print(f"backup → {bak}")

    prs = Presentation(FILE)
    slide = prs.slides[7]  # 第8页

    # 删掉正文占位符（留着会显示空框），以及旧图/旧文本框
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            if shape.placeholder_format.idx == 1:  # body
                shape._element.getparent().remove(shape._element)
            continue
        if shape.shape_type == 13 or shape.has_text_frame:
            shape._element.getparent().remove(shape._element)

    # 标题拉回标准左缘（原先偏右）
    title = next(sh for sh in slide.placeholders
                 if sh.placeholder_format.idx == 0)
    title.left, title.width = Inches(0.53), Inches(9.49)
    # 确保标题有 noProof
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            style_run(run, Pt(26), bold=True)

    slide.shapes.add_picture(IMG, Inches(IMG_LEFT), Inches(IMG_TOP),
                             width=Inches(IMG_WIDTH))
    disable_proofing(slide)
    prs.save(FILE)
    print(f"slide 8 updated → {FILE}")


if __name__ == "__main__":
    main()
