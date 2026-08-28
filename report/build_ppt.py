# -*- coding: utf-8 -*-
"""
根据「实习小结模板.pptx」生成「实习小结_毛泽辉.pptx」。

内容来源：实习小结_方案讨论.md（已定稿的次要部分）。
主体四页（工作目标及完成情况）目前留空，待内容确定后填入 MAIN_TITLES 下方重跑即可。

用法：python build_ppt.py
"""

import re

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

TEMPLATE = "实习小结模板.pptx"
OUTPUT = "实习小结_毛泽辉.pptx"

PLACEHOLDER_TEXT = "（内容待补充）"

EA_FONT = "微软雅黑"      # 中文字形
LATIN_FONT = "Segoe UI"   # 西文与数字字形

TITLE_SIZE = Pt(26)
BODY_SIZE = Pt(18)        # 正文一级
SUB_SIZE = Pt(15)         # 正文二级


# ---------------------------------------------------------------- 字体与排版

def style_run(run, size, bold=None):
    """统一中西文字体、字号，并关闭拼写检查（消除红色波浪线）。"""
    font = run.font
    font.size = size
    if bold is not None:
        font.bold = bold
    font.name = LATIN_FONT              # 写入 a:latin

    rPr = font._rPr
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        latin = rPr.find(qn("a:latin"))
        latin.addnext(ea) if latin is not None else rPr.append(ea)
    ea.set("typeface", EA_FONT)

    rPr.set("lang", "zh-CN")
    rPr.set("altLang", "en-US")
    rPr.set("noProof", "1")             # 关键：不做拼写校对


def disable_proofing(prs):
    """兜底：对未经本脚本写入的文本（封面、Q&A、结尾页）同样关闭拼写检查。"""
    def walk(shapes):
        for shape in shapes:
            if shape.shape_type == 6:   # GROUP
                walk(shape.shapes)
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run.font._rPr
                    rPr.set("lang", "zh-CN")
                    rPr.set("altLang", "en-US")
                    rPr.set("noProof", "1")

    for slide in prs.slides:
        walk(slide.shapes)


def reset_bullet(para):
    """清掉段落自带的项目符号与缩进覆盖，改为继承版式。

    text_frame.clear() 会保留首段的段落属性；模板某些页首行原本是无符号的
    引导语（如「包括但不限于」），不清掉会导致首项缺项目符号、缩进也不一致。
    """
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    for attr in ("marL", "indent"):
        pPr.attrib.pop(attr, None)


def set_rich_text(text_frame, items, body_size=BODY_SIZE, sub_size=SUB_SIZE,
                  line_spacing=1.2, space_after=Pt(10)):
    """把 [(层级, 文本), ...] 写进占位符。文本中 **xx** 渲染为加粗。"""
    text_frame.clear()
    text_frame.word_wrap = True

    for i, (level, raw) in enumerate(items):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.level = level
        reset_bullet(para)
        para.line_spacing = line_spacing
        para.space_before = Pt(0)
        para.space_after = space_after
        size = body_size if level == 0 else sub_size

        for seg in re.split(r"(\*\*.+?\*\*)", raw):
            if not seg:
                continue
            run = para.add_run()
            bold = seg.startswith("**") and seg.endswith("**")
            run.text = seg[2:-2] if bold else seg
            style_run(run, size, bold=True if bold else None)


def set_title(slide, text):
    tf = get_ph(slide, 0).text_frame
    tf.text = text
    para = tf.paragraphs[0]
    para.line_spacing = 1.1
    for run in para.runs:
        style_run(run, TITLE_SIZE, bold=True)


def get_ph(slide, idx):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    raise KeyError(f"占位符 idx={idx} 不存在")


def fill(slide, title, items, **kwargs):
    set_title(slide, title)
    set_rich_text(get_ph(slide, 1).text_frame, items, **kwargs)


def add_content_slide(prs, ref_slide, title, items, **kwargs):
    """按 ref_slide 的版式与占位符位置新建一页。"""
    slide = prs.slides.add_slide(ref_slide.slide_layout)

    # 删除版式带来的、本模板不使用的多余占位符（保留标题/正文/页码）
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx not in (0, 1, 4):
            shape._element.getparent().remove(shape._element)

    # 位置与尺寸对齐参考页，避免新页与原有页错位
    for idx in (0, 1, 4):
        try:
            src, dst = get_ph(ref_slide, idx), get_ph(slide, idx)
        except KeyError:
            continue
        dst.left, dst.top, dst.width, dst.height = src.left, src.top, src.width, src.height

    fill(slide, title, items, **kwargs)
    return slide


def reorder(prs, order):
    """按给定的原始索引顺序重排幻灯片。"""
    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    for entry in entries:
        id_list.remove(entry)
    for i in order:
        id_list.append(entries[i])


# ---------------------------------------------------------------- 页面内容

COVER_LINE = "AI Lab - 毛泽辉\n时间：8.20-8.31"   # TODO 实习起止时间待确认

OUTLINE = [
    (0, "自我介绍"),
    (0, "工作目标及完成情况"),
    (0, "实习体会"),
    (0, "职业规划"),
]

SELF_INTRO = [
    (0, "教育背景"),
    (1, "本科：浙江大学 物理学系"),
    (1, "硕士：南洋理工大学 应用人工智能"),
    (0, "专业方向：通用人工智能"),
    (0, "目前职位：通用人工智能研究　｜　导师：张亚红"),
    (0, "个人爱好：游戏、动漫"),
]

MAIN_TITLES = [
    "工作目标及完成情况（一）｜自进化调研",
    "工作目标及完成情况（二）｜OpenClaw-RL 实现方法",
    "工作目标及完成情况（三）｜复现实验结果",
    "工作目标及完成情况（四）｜方法泛化方案",
]

# 两条收获合并为一页：一级为收获标题，二级为要点
GAINS = [
    (0, "**收获一 ｜ 对 Agent 自进化形成了系统认知**"),
    (1, "**结论**：自进化是通往 AGI 的必要路径；当前没有一条技术路线是完备的，值得多方案并行探索"),
    (1, "**必要性**：模型上线即冻结，而用户任务分布持续漂移 —— 能力衰减是必然"),
    (1, "**技术谱系**：**非参数化**（记忆 / 技能库；零停机、见效快，但模型能力本身不变）"
        "　↔　**参数化**（权重更新；真正长能力，但代价高、传统需停机）"),
    (1, "**我的落点**：**复现并验证了 OpenClaw-RL 的核心机制** —— 属于参数化一端，"
        "且把权重更新做进真实对话流，实现“边服务边进化”"),
    (0, "**收获二 ｜ 试用 multica，体验了另一条技术路线**"),
    (1, "试用了自研多智能体协作平台 multica，并在其上实际完成了一个完整产品的开发"),
    (1, "体验到**非参数化自进化**的实际形态 —— 与收获一的参数化路线互为对照，"
        "印证了“两条路互补”的判断"),
    (1, "直观看到了多智能体协作的优势；结合实际使用体验，向平台提出了多项改进意见"),
]

IMPRESSION = [
    (0, "作为 **AI 发起者小组组长**，与其他部门的实习生组队协作"),
    (0, "五周时间，从 0 到 1 完成了「掉了么」这个可视化产品"),
    (0, "也看到了其他小组的成果 —— 是一次难得的横向视野"),
]

CAREER = [
    (0, "**职业目标**：成为能在所处方向上做出行业领先成果的研究员 —— **从跟随前沿，到定义前沿**"),
    (0, "**短期**：从“复现”走向“独立” —— 能自己判断方向价值、设计实验路径，做出经得起推敲的结果"),
    (0, "**中期**：做出**行业领先的成果** —— 在具体问题上达到 SOTA，并推动其在真实场景中落地"),
    (0, "**长期**：提出**有原创性的方法或理论**，在所处方向上形成影响力，让自己的工作成为别人的参照"),
    (0, "**行动方案**：在自进化主线上保持高频跟进；把复现能力沉淀为可复用的实验基建，"
        "缩短从论文到验证的周期；成果尽量对外呈现，接受同行检验"),
]


# ---------------------------------------------------------------- 构建

def main():
    prs = Presentation(TEMPLATE)

    # 模板原始页序：0封面 1大纲 2自我介绍 3工作目标 4实习体会 5职业规划 6Q&A 7结尾
    cover, outline, intro = prs.slides[0], prs.slides[1], prs.slides[2]
    main_1, gains, career = prs.slides[3], prs.slides[4], prs.slides[5]

    # 封面：逐行替换，保留模板原有字号与位置
    for shape in cover.shapes:
        if shape.has_text_frame and "Lab" in shape.text_frame.text:
            lines = COVER_LINE.split("\n")
            for i, para in enumerate(shape.text_frame.paragraphs):
                if i < len(lines) and para.runs:
                    para.runs[0].text = lines[i]
                    for extra in para.runs[1:]:
                        extra.text = ""

    # 字号按每页信息密度分别设定：条目少的页放大，密的页收紧，使各页饱满度接近
    fill(outline, "呈现大纲", OUTLINE, body_size=Pt(24), space_after=Pt(22))
    fill(intro, "自我介绍", SELF_INTRO,
         body_size=Pt(22), sub_size=Pt(18), space_after=Pt(16))

    # 主体第一页复用模板原「工作目标及完成情况」页，其余三页新建
    fill(main_1, MAIN_TITLES[0], [(0, PLACEHOLDER_TEXT)], body_size=Pt(20))
    for title in MAIN_TITLES[1:]:
        add_content_slide(prs, main_1, title, [(0, PLACEHOLDER_TEXT)], body_size=Pt(20))

    # 两条收获合并一页（本页最密，字号最小）
    fill(gains, "实习收获", GAINS,
         body_size=Pt(18), sub_size=Pt(15), line_spacing=1.15, space_after=Pt(9))
    add_content_slide(prs, gains, "印象最深的一件事 ｜ 联想暑期活动", IMPRESSION,
                      body_size=Pt(24), space_after=Pt(28))

    fill(career, "职业规划", CAREER, space_after=Pt(14))

    # 新建页追加在末尾（8,9,10=工作②③④　11=印象最深），重排回正确顺序
    reorder(prs, [0, 1, 2, 3, 8, 9, 10, 4, 11, 5, 6, 7])

    disable_proofing(prs)
    prs.save(OUTPUT)
    print(f"已生成 {OUTPUT}，共 {len(prs.slides._sldIdLst)} 页")


if __name__ == "__main__":
    main()
