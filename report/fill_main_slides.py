# -*- coding: utf-8 -*-
"""填充「工作目标及完成情况」主体四页的内容（就地修改，不重建文件）。

目前已实现：第4页 自进化调研（左文右图两栏）。
第5~7页待内容定稿后在下方 fill_slide5/6/7 里补。

可重复运行：每次会先清掉该页原有图片再重新插入，不会叠加。

文案来源：实习小结_方案讨论.md「第4页 自进化调研」节。

用法：python fill_main_slides.py
"""

import re

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FILE = "实习小结_毛泽辉.pptx"

EA_FONT = "微软雅黑"
LATIN_FONT = "Segoe UI"

# --- 第4页版式（单位：英寸）------------------------------------------------
TEXT_LEFT, TEXT_WIDTH = 0.53, 5.80          # 左栏文字，与标题同基准线
TEXT_TOP, TEXT_HEIGHT = 1.75, 5.00      # 底部留出参考文献行的位置
REF4 = "[1] Self-Improving Agents in the Era of Experience: A Survey of Self- to Meta-Evolution (2026)"
REF4_TOP = 6.92

IMG_PATH = "personal_agent_self_evolution_concept.png"
IMG_LEFT, IMG_WIDTH = 6.55, 6.45            # 右栏配图；原图长宽比 1.769
IMG_TOP = 2.55                              # 高度约 3.65，在正文区内垂直居中

HEAD_SIZE, BODY_SIZE = Pt(16), Pt(14)

# 主体四页标题。不沿用模板的「工作目标及完成情况」，也不编序号——
# 四页统一为「工作完成情况」，不单设工作目标页。
MAIN_TITLES = [
    "工作完成情况 ｜ 自进化调研",
    "工作完成情况 ｜ OpenClaw-RL 实现方法",
    "工作完成情况 ｜ 复现实验结果",
    "工作完成情况 ｜ 方法泛化方案",
]
OUTLINE_ITEM_OLD, OUTLINE_ITEM_NEW = "工作目标及完成情况", "工作完成情况"

# 模板标题占位符的标准位置，各页保持一致
TITLE_LEFT, TITLE_WIDTH = 0.53, 9.49

# --- 第9页版式（印象最深的一件事：文字在上，三张 4:3 照片排一行）-------------
TEXT9_LEFT, TEXT9_TOP = 0.53, 1.62
TEXT9_WIDTH, TEXT9_HEIGHT = 12.27, 2.30
PHOTO9_W, PHOTO9_GAP, PHOTO9_TOP = 3.83, 0.32, 4.10
PHOTO9_RATIO = 4 / 3

# --- 第5页版式 -------------------------------------------------------------
# 细节图是竖版（1840×2600，长宽比 0.708），占右侧整幅高度，
# 标题与正文一并收进左栏。
IMG5_HEIGHT = 6.80          # 略小于满幅，四周留呼吸空间
IMG5_RIGHT = 13.10          # 右边缘位置
TEXT5_LEFT, TEXT5_WIDTH = 0.53, 7.50

# 左栏不复述图，而是给一把读图的钥匙：四组依次对应图上从上到下四个区域
SLIDE5 = [
    (0, "**① 三个角色，一套真实环境**"),
    (1, "Simulator（Qwen3-32B）按人设扮演 Student / TA / Teacher，"
        "在 OpenClaw 里与 Policy（Qwen3-4B）对话"),
    (0, "**② 两种训练顺序**"),
    (1, "**Separate**（本次已完成）：三角色各训一个模型，须依次进行"
        "　｜　**Joint**：三角色共享一个模型，同时训练"),
    (0, "**③ 一次回复，两条去向**"),
    (1, "对话循环：回复照常送回用户　｜　训练循环：同一次回复交给判官"),
    (0, "**④ 两类信号，同一个判官**"),
    (1, "Eval 判官答“好不好”→ GRPO；Hint 判官答“该怎么改”→ OPD。"
        "同一份 PRM 权重，只是提示词不同"),
    (0, "**⑤ 一次反传，权重回流**"),
    (1, "两个损失加权合并后一次更新，新权重同步回线上 Policy"),
]
# 落点句：无项目符号，与上方四组区分开
SLIDE5_CLOSING = "训练与对话**异步并行**：回应照常送出，权重在线更新，全程不停机"

# --- 第6页版式 -------------------------------------------------------------
# 左：论文 Table 3 原图（直接裁自 PDF，仅高亮 19.2 一格）；右：对话对照图
TBL6_PATH = "table3_highlighted.png"         # 长宽比 3.48（很扁）
TBL6_LEFT, TBL6_TOP, TBL6_WIDTH = 0.53, 1.62, 5.55   # 与右侧图顶部对齐

IMG6_PATH = "dialogue_evidence.png"          # 长宽比 1.232
IMG6_LEFT, IMG6_TOP, IMG6_HEIGHT = 6.10, 1.60, 5.60

CAP6 = "论文 Table 3 · 达到收敛所需的 session 数（越小越好）"
NOTE6 = "高亮格 19.2 = 本次复现对应的 Separate / Student / Hybrid RL"

# 左下说明块：把「论文 19.2 vs 本次 20」这个核心结论直接摆出来，
# 不让评委自己去左右两栏之间拼。
RESULT6_TITLE = "本次复现结果"
RESULT6_ROWS = [                      # (标签, 数字)
    ("论文（5 次独立试验均值）", "19.2"),
    ("本次复现（单次运行）", "20"),
]
RESULT6_NOTES = [
    # 配置行较长，手动断成两行，避免自动折行把模型名从中间劈开
    "配置：Policy / PRM = Qwen3-4B-Thinking-2507",
    "　　　Simulator = DeepSeek-V4-Flash　｜　GSM8K　｜　Separate-Student",
    "判定：rule-based，连续 3 个 session 的 Turn 1 回复均满足偏好",
    # 分母 15 已剔除第 32、34 个 session——那两次是「无回复 / 生成失败」，
    # 因不含 bold 与列表会被格式规则误判为达标
    "第 18 个 session 起共 15 次有效回复，14 次 Turn 1 直接达标，仅 1 次退化",
]

# --- 第7页版式（改版：右栏放迁移对照图，左栏改成"缩小的 Table 1 对比图 + 文字"）---
# 「怎么迁」整段文字改由右侧迁移对照图承担，不再用文字重复一遍。

# 右栏：迁移对照图（复现细节图基础上标注改/不改/新增），长宽比 2400/2600≈0.923
DIAGRAM7_PATH = "metaclaw_migration_diagram.png"
DIAGRAM7_LEFT, DIAGRAM7_TOP, DIAGRAM7_HEIGHT = 7.55, 1.55, 5.75

# 左上：Table 1 对比图缩小后放置，长宽比 1.511（"缩小一点点"，非大改）
IMG7_PATH = "table1_with_ours.png"
IMG7_LEFT, IMG7_TOP, IMG7_WIDTH = 0.53, 1.62, 4.35

NOTE7 = ("本次仅迁移 RL 训练方法，不含 MetaClaw 技能库；模型与评测环境均不同，"
         "与论文各行为数量级参照")

# 左下：文字缩到 ①③④ 三段，②「怎么迁」整段删除（右图已经讲清楚）
TEXT7_LEFT, TEXT7_WIDTH = 0.53, 4.85
# TEXT7_TOP 由 fill_slide7() 按 Table 1 图实际高度动态算出，此处仅占位
TEXT7_TOP, TEXT7_HEIGHT = 5.35, 1.95

# 最终版：只剩①②两段，②不提 bench 复核（论文 Table 1 Full 本身就是 live
# 聚合，我们跟论文方法学一致，不是需要补齐的偏差；bench 复核只是论文没做过
# 的额外校验，不是当务之急，用户决定删掉更聚焦）。
# 「怎么迁」「结果」两段已删——内容改由右侧 metaclaw_migration_diagram.png
# 和左上 table1_with_ours.png 承担，不再用文字重复。
SLIDE7 = [
    (0, "**① 任务背景**"),
    (1, "方法已在 Personal Agent 场景复现校准，换到 **MetaClaw-Bench**"
        "（另一篇论文的评测基准，30 天连续任务流）验证普适性"),
    (0, "**② 未来提升方向**"),
    (1, "目标：训练能稳定、有效地跑满完整 30 天"),
    (1, "已定位并修复影响训练后期效率的根因，下一步在真实训练中验证效果"),
]

# ①② 改用综述原文并标出处（此前的转述"方向太大、又没有来源"，立不住）；
# ③ 改为本次实验的第一手体会，不引文献。
SLIDE4 = [
    (0, "**① 自进化的理解**　（综述 [1]）"),
    (1, "“…agentic AI is no longer defined only by what a model can infer from "
        "static data, but by how a deployed system accumulates, organizes, and "
        "reuses experience from interaction.”"),
    (0, "**② 经验的五个更新层面**　（同上）"),
    (1, "reusable skill / persistent memory / verifiable environment feedback / "
        "**trainable model behavior** / meta-level control"),
    (0, "**③ 本次实验的主要难点**"),
    (1, "不在 RL 算法本身，而在“经验 → 更新”这条链路：信号怎么提取、"
        "怎么区分环境故障与模型失败、训练怎么在真实运行时里跑"),
    (0, "**④ 本次工作的定位**"),
    (1, "**OpenClaw-RL** ＝ **模型权重**这一层，训练直接跑在真实对话流中"
        " —— 对话本身就是训练信号"),
]


def style_run(run, size, bold=None):
    """统一中西文字体与字号，并关闭拼写检查（否则中文显示红色波浪线）。"""
    font = run.font
    font.size = size
    if bold is not None:
        font.bold = bold
    font.name = LATIN_FONT

    rPr = font._rPr
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        latin = rPr.find(qn("a:latin"))
        latin.addnext(ea) if latin is not None else rPr.append(ea)
    ea.set("typeface", EA_FONT)

    rPr.set("lang", "zh-CN")
    rPr.set("altLang", "en-US")
    rPr.set("noProof", "1")


def reset_bullet(para):
    """清掉段落自带的项目符号与缩进覆盖，改为继承版式。"""
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    for attr in ("marL", "indent"):
        pPr.attrib.pop(attr, None)


def set_no_bullet(para):
    """去掉该段的项目符号（用于收尾句，使其与上方分组区分开）。"""
    pPr = para._p.get_or_add_pPr()
    pPr.attrib["marL"] = "0"
    pPr.attrib["indent"] = "0"
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def set_rich_text(text_frame, items, head_size, body_size,
                  line_spacing=1.15, space_after=Pt(8)):
    """写入 [(层级, 文本), ...]；文本中 **xx** 渲染为加粗。"""
    text_frame.clear()
    text_frame.word_wrap = True

    for i, (level, raw) in enumerate(items):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.level = level
        reset_bullet(para)
        para.line_spacing = line_spacing
        para.space_before = Pt(0)
        # 每组小标题前留出间隔，让四组在视觉上分开
        para.space_after = space_after if level else Pt(4)
        size = head_size if level == 0 else body_size

        for seg in re.split(r"(\*\*.+?\*\*)", raw):
            if not seg:
                continue
            run = para.add_run()
            bold = seg.startswith("**") and seg.endswith("**")
            run.text = seg[2:-2] if bold else seg
            style_run(run, size, bold=True if bold else None)


def get_body(slide):
    return next(sh for sh in slide.placeholders
                if sh.placeholder_format.idx == 1)


def fill_slide4(slide):
    body = get_body(slide)
    body.left, body.width = Inches(TEXT_LEFT), Inches(TEXT_WIDTH)
    body.top, body.height = Inches(TEXT_TOP), Inches(TEXT_HEIGHT)
    # 加了引文与页脚后条目变长，字号较初版收紧一档
    set_rich_text(body.text_frame, SLIDE4, Pt(15), Pt(13),
                  line_spacing=1.12, space_after=Pt(7))

    # 先清掉旧图与旧页脚，保证脚本可重复运行
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        if shape.shape_type == 13 or shape.has_text_frame:
            shape._element.getparent().remove(shape._element)

    # 页脚：完整参考文献。标题行只标角标 [1]，避免 34 词的全称把引文挤下去
    ref = slide.shapes.add_textbox(Inches(TEXT_LEFT), Inches(REF4_TOP),
                                   Inches(TEXT_WIDTH), Inches(0.42))
    ref.text_frame.word_wrap = True
    para = ref.text_frame.paragraphs[0]
    para.line_spacing = 1.12
    run = para.add_run()
    run.text = REF4
    style_run(run, Pt(9.5))

    slide.shapes.add_picture(IMG_PATH, Inches(IMG_LEFT), Inches(IMG_TOP),
                             width=Inches(IMG_WIDTH))


def fill_slide7(slide):
    """右栏：迁移对照图（大，接近满高）。左上：缩小的 Table 1 对比图。
    左下：①③④ 三段文字（②「怎么迁」整段删除，改由右图承担）。
    """
    # 清场：删掉上一次运行生成的图片与说明文字，保证脚本可重复运行
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        if shape.shape_type == 13 or shape.has_text_frame:
            shape._element.getparent().remove(shape._element)

    # 右栏：迁移对照图
    slide.shapes.add_picture(DIAGRAM7_PATH, Inches(DIAGRAM7_LEFT), Inches(DIAGRAM7_TOP),
                             height=Inches(DIAGRAM7_HEIGHT))

    # 左上：Table 1 对比图（缩小版）
    slide.shapes.add_picture(IMG7_PATH, Inches(IMG7_LEFT), Inches(IMG7_TOP),
                             width=Inches(IMG7_WIDTH))
    img_h = IMG7_WIDTH / 1.511

    note_top = IMG7_TOP + img_h + 0.08
    note = slide.shapes.add_textbox(Inches(IMG7_LEFT), Inches(note_top),
                                    Inches(IMG7_WIDTH), Inches(0.5))
    note.text_frame.word_wrap = True
    para = note.text_frame.paragraphs[0]
    para.line_spacing = 1.08
    run = para.add_run()
    run.text = NOTE7
    style_run(run, Pt(10))

    # 左下：文字区紧跟在说明小字下方，动态起算（避免跟图/小字重叠）
    text_top = note_top + 0.38
    body = get_body(slide)
    body.left, body.width = Inches(TEXT7_LEFT), Inches(TEXT7_WIDTH)
    body.top, body.height = Inches(text_top), Inches(7.32 - text_top)
    set_rich_text(body.text_frame, SLIDE7, Pt(13.5), Pt(11.5),
                  line_spacing=1.03, space_after=Pt(4))


def fill_slide6(slide):
    """左：论文 Table 3 原图（高亮 19.2）；右：本次复现的对话对照图。"""
    # 本页只有图和说明文字，正文占位框整个删掉——留着会在编辑态显示
    # 「单击此处添加文本」的空框
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx == 1:
            shape._element.getparent().remove(shape._element)

    # 清场：删掉上一次运行生成的图片与说明文字，保证脚本可重复运行
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        if shape.has_table or shape.shape_type == 13 or shape.has_text_frame:
            shape._element.getparent().remove(shape._element)

    def caption(top, text, size, bold):
        box = slide.shapes.add_textbox(Inches(TBL6_LEFT), Inches(top),
                                       Inches(TBL6_WIDTH), Inches(0.32))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = text
        style_run(run, Pt(size), bold=bold)

    caption(TBL6_TOP, CAP6, 13, True)
    slide.shapes.add_picture(TBL6_PATH, Inches(TBL6_LEFT), Inches(TBL6_TOP + 0.42),
                             width=Inches(TBL6_WIDTH))
    # 表图高度 = 5.55 / 3.48 ≈ 1.60
    caption(TBL6_TOP + 0.42 + 1.60 + 0.14, NOTE6, 11, False)

    result_block(slide)

    slide.shapes.add_picture(IMG6_PATH, Inches(IMG6_LEFT), Inches(IMG6_TOP),
                             height=Inches(IMG6_HEIGHT))


def result_block(slide, top=4.30):
    """左下角结论块：论文值与本次复现值直接并列，下方补口径与配置。"""
    box = slide.shapes.add_textbox(Inches(TBL6_LEFT), Inches(top),
                                   Inches(TBL6_WIDTH), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True

    def para(first=False):
        return tf.paragraphs[0] if first else tf.add_paragraph()

    p0 = para(True)
    p0.space_after = Pt(8)
    style_run(add(p0, RESULT6_TITLE), Pt(15), bold=True)

    for label, value in RESULT6_ROWS:
        p = para()
        p.space_after = Pt(2)
        style_run(add(p, f"{label}　"), Pt(12))
        style_run(add(p, value), Pt(24), bold=True)

    for i, note in enumerate(RESULT6_NOTES):
        p = para()
        p.space_before = Pt(10) if i == 0 else Pt(1 if i == 1 else 4)
        p.space_after = Pt(2)
        p.line_spacing = 1.15
        style_run(add(p, note), Pt(10.5))


def add(para, text):
    run = para.add_run()
    run.text = text
    return run


def layout_slide9(slide):
    """第9页（印象最深的一件事）：文字在上，三张照片等大排成一行在下。

    照片是用户手工插入的，这里只做等比缩放与定位，不删除、不替换。
    原摆法是「一张右上 + 两张底部」，右上那张会压住第一条文字，且三张尺寸不一致。
    """
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    if len(pics) != 3:
        raise RuntimeError(f"第9页图片数量为 {len(pics)}，与预期的 3 张不符，已中止")

    # 按时间顺序重排：工作 → 路演 → 合影领奖，结尾落在领奖照上。
    # 合影原本单独放在上方，用 top 区分；另两张按原本的左右顺序。
    top_one = min(pics, key=lambda sh: sh.top)
    bottom_two = sorted((sh for sh in pics if sh is not top_one),
                        key=lambda sh: sh.left)
    ordered = bottom_two + [top_one]

    title = next(sh for sh in slide.placeholders
                 if sh.placeholder_format.idx == 0)
    title.left, title.width = Inches(TITLE_LEFT), Inches(TITLE_WIDTH)

    body = get_body(slide)
    body.left, body.top = Inches(TEXT9_LEFT), Inches(TEXT9_TOP)
    body.width, body.height = Inches(TEXT9_WIDTH), Inches(TEXT9_HEIGHT)

    span = PHOTO9_W * 3 + PHOTO9_GAP * 2
    left = (13.333 - span) / 2                    # 整排照片在页面上水平居中
    for pic in ordered:
        pic.width = Inches(PHOTO9_W)
        pic.height = Inches(PHOTO9_W / PHOTO9_RATIO)
        pic.left, pic.top = Inches(left), Inches(PHOTO9_TOP)
        left += PHOTO9_W + PHOTO9_GAP


def layout_slide5(slide):
    """只调版式，不动内容：缩小右侧竖版细节图，标题与正文收进左栏。

    图片是用户手工插入的，这里只做等比缩放与定位，**不删除、不替换**。
    """
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    if len(pics) != 1:
        raise RuntimeError(f"第5页图片数量为 {len(pics)}，与预期的 1 张不符，已中止")

    pic = pics[0]
    ratio = pic.width / pic.height
    pic.height = Inches(IMG5_HEIGHT)
    pic.width = Inches(IMG5_HEIGHT * ratio)
    pic.left = Inches(IMG5_RIGHT - IMG5_HEIGHT * ratio)
    pic.top = Inches((7.5 - IMG5_HEIGHT) / 2)

    title = next(sh for sh in slide.placeholders
                 if sh.placeholder_format.idx == 0)
    title.left, title.width = Inches(TEXT5_LEFT), Inches(TEXT5_WIDTH)

    body = get_body(slide)
    body.left, body.width = Inches(TEXT5_LEFT), Inches(TEXT5_WIDTH)
    set_rich_text(body.text_frame, SLIDE5 + [(0, SLIDE5_CLOSING)],
                  HEAD_SIZE, BODY_SIZE)
    closing = body.text_frame.paragraphs[-1]
    closing.space_before = Pt(14)
    set_no_bullet(closing)


def retitle_main_slides(prs):
    """四页统一标题，并同步更新大纲页对应条目。"""
    main_slides = [prs.slides[i] for i in range(3, 7)]   # slides 不支持切片
    for slide, title in zip(main_slides, MAIN_TITLES):
        ph = next(sh for sh in slide.placeholders
                  if sh.placeholder_format.idx == 0)
        para = ph.text_frame.paragraphs[0]
        if not para.runs:
            raise RuntimeError(f"标题占位符为空：{title}")
        para.runs[0].text = title
        for extra in para.runs[1:]:
            extra._r.getparent().remove(extra._r)

    for para in get_body(prs.slides[1]).text_frame.paragraphs:
        for run in para.runs:
            if OUTLINE_ITEM_OLD in run.text:
                run.text = run.text.replace(OUTLINE_ITEM_OLD, OUTLINE_ITEM_NEW)


def main():
    prs = Presentation(FILE)
    retitle_main_slides(prs)
    fill_slide4(prs.slides[3])
    layout_slide5(prs.slides[4])
    fill_slide6(prs.slides[5])
    fill_slide7(prs.slides[6])
    layout_slide9(prs.slides[8])
    prs.save(FILE)
    print(f"已更新四页标题；第4~7页已填充；第9页照片排版已调整 → {FILE}")


if __name__ == "__main__":
    main()
