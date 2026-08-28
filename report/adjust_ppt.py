# -*- coding: utf-8 -*-
"""对已有的 实习小结_毛泽辉.pptx 做安全的就地维护（不改版式、不动图片）。

只做两件事：
  1. 补上缺失的页码占位符
  2. 全篇补齐 noProof —— 在 PowerPoint 里手工输入的文字不带此属性，
     中文会被按英文拼写校对而显示红色波浪线。每次手工改完文字后重跑本脚本即可。

⚠️ 本脚本刻意不包含任何版式/图片/文字内容调整：第3页的校徽位置与尺寸、
   封面日期等都是用户手动设定的，脚本不得覆盖。
   （2026-08-21 教训：曾经在这里硬编码封面日期并强制重写，导致用户手动
   改的日期被这个脚本悄悄覆盖回去。已删除该逻辑——封面内容只手动改，
   不再由脚本管理。）

用法：python adjust_ppt.py
"""

import copy

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

FILE = "实习小结_毛泽辉.pptx"


def ensure_page_numbers(prs):
    """补上缺失的页码占位符。

    脚本新建的页面不会自动带上版式里的页码框，导致部分内容页有页码、
    部分没有。这里从已有页码的页复制一份过去。封面与结尾页用 Blank 版式、
    本来就没有占位符，跳过。
    """
    def page_nums(slide):
        # 按占位符类型判断，不能按 idx——已有页码框的 idx 表示方式并不统一
        return [sh for sh in slide.placeholders
                if sh.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER]

    def has_title(slide):
        return any(sh.placeholder_format.idx == 0 for sh in slide.placeholders)

    donor = next((s for s in prs.slides if page_nums(s)), None)
    if donor is None:
        raise RuntimeError("没有任何页面带页码占位符，无法复制")
    src = page_nums(donor)[0]

    fixed = []
    for i, slide in enumerate(prs.slides, 1):
        existing = page_nums(slide)
        for extra in existing[1:]:               # 清掉重复的页码框
            extra._element.getparent().remove(extra._element)
        if has_title(slide) and not existing:    # 是内容页但没有页码
            slide.shapes._spTree.append(copy.deepcopy(src._element))
            fixed.append(i)
    return fixed


def disable_proofing(prs):
    count = 0
    def walk(shapes):
        nonlocal count
        for shape in shapes:
            if shape.shape_type == 6:      # GROUP
                walk(shape.shapes)
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run.font._rPr
                    if rPr.get("noProof") != "1":
                        count += 1
                    rPr.set("lang", "zh-CN")
                    rPr.set("altLang", "en-US")
                    rPr.set("noProof", "1")

    for slide in prs.slides:
        walk(slide.shapes)
    return count


def main():
    prs = Presentation(FILE)
    pages = ensure_page_numbers(prs)
    fixed = disable_proofing(prs)
    prs.save(FILE)
    print(f"补页码：{pages or '无'}；补 noProof：{fixed} 处")


if __name__ == "__main__":
    main()
