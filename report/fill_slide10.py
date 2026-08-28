# -*- coding: utf-8 -*-
"""只更新第10页：清空旧正文与旧图，插入 career_plan.png。

用法：python fill_slide10.py
"""

import os
import shutil
from datetime import datetime

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

FILE = "实习小结_毛泽辉.pptx"
IMG = "career_plan.png"

EA_FONT = "微软雅黑"
LATIN_FONT = "Segoe UI"

# 标题下方内容区；图 1952×912 ≈ 2.14
IMG_LEFT, IMG_TOP = 0.55, 1.48
IMG_WIDTH = 12.20


def style_run(run, size, bold=None):
    font = run.font
    font.size = size
    if bold is not None:
        font.bold = bold
    font.name = LATIN_FONT
    rPr = font._rPr
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        latin = rPr.find(qn("a:latin"))
        latin.addnext(ea) if latin is not None else rPr.append(ea)
    ea.set("typeface", EA_FONT)
    rPr.set("lang", "zh-CN")
    rPr.set("altLang", "en-US")
    rPr.set("noProof", "1")


def disable_proofing(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rPr = run.font._rPr
                rPr.set("lang", "zh-CN")
                rPr.set("altLang", "en-US")
                rPr.set("noProof", "1")


def main():
    stamp = datetime.now().strftime("%H%M")
    bak = f"实习小结_毛泽辉_backup_改第10页前_{stamp}.pptx"
    shutil.copy2(FILE, bak)
    print(f"backup → {bak}")

    prs = Presentation(FILE)
    slide = prs.slides[9]  # 第10页

    # 删正文占位符 + 非标题/页码的图与文本框（含满幅旧图）
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            if idx == 1:  # body
                shape._element.getparent().remove(shape._element)
            continue
        if shape.shape_type == 13 or shape.has_text_frame:
            shape._element.getparent().remove(shape._element)

    title = next(sh for sh in slide.placeholders
                 if sh.placeholder_format.idx == 0)
    # 与模板一致：标题框在 left≈1.92（勿改成 0.53，会偏左）
    title.left, title.top = Inches(1.919), Inches(0.317)
    title.width, title.height = Inches(9.495), Inches(1.318)
    # 确保标题文字正确
    para = title.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = "职业规划"
        for extra in para.runs[1:]:
            extra._r.getparent().remove(extra._r)
        style_run(para.runs[0], Pt(26), bold=True)
    else:
        run = para.add_run()
        run.text = "职业规划"
        style_run(run, Pt(26), bold=True)

    slide.shapes.add_picture(IMG, Inches(IMG_LEFT), Inches(IMG_TOP),
                             width=Inches(IMG_WIDTH))
    disable_proofing(slide)
    # 若成品被占用，先写临时文件再替换
    tmp = FILE.replace(".pptx", "._tmp_slide10.pptx")
    prs.save(tmp)
    try:
        os.replace(tmp, FILE)
        print(f"slide 10 updated → {FILE}")
    except PermissionError:
        print(f"成品被占用，已保存到 → {tmp}")
        print("请关闭 PowerPoint 中的该文件后重跑：python fill_slide10.py")
        raise


if __name__ == "__main__":
    main()
